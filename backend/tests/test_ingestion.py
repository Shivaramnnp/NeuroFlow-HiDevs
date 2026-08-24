from __future__ import annotations

import io
import json
import uuid
import pytest
from unittest.mock import AsyncMock, MagicMock, patch

from PIL import Image
import docx
import pandas as pd
from pptx import Presentation

from backend.pipelines.ingestion.extractors.base import ExtractedPage
from backend.pipelines.ingestion.extractors.pdf_extractor import PDFExtractor, format_markdown_table
from backend.pipelines.ingestion.extractors.docx_extractor import DocxExtractor
from backend.pipelines.ingestion.extractors.image_extractor import ImageExtractor
from backend.pipelines.ingestion.extractors.csv_extractor import CSVExtractor
from backend.pipelines.ingestion.extractors.url_extractor import URLExtractor
from backend.pipelines.ingestion.extractors.pptx_extractor import PPTXExtractor
from backend.pipelines.ingestion.chunker import (
    Chunk,
    count_tokens,
    fixed_size_chunking,
    semantic_chunking,
    hierarchical_chunking,
    select_chunking_strategy,
    chunk_pages,
)
from backend.pipelines.ingestion.pipeline import (
    compute_content_hash,
    check_deduplication,
    IngestionPipeline,
)
from backend.providers.base import GenerationResult
from backend.providers.client import NeuroFlowClient


# --- 1. Deduplication Tests ---

def test_compute_content_hash():
    data1 = b"NeuroFlow test content"
    data2 = b"NeuroFlow test content"
    data3 = b"Different content"

    hash1 = compute_content_hash(data1)
    hash2 = compute_content_hash(data2)
    hash3 = compute_content_hash(data3)

    assert len(hash1) == 64
    assert hash1 == hash2
    assert hash1 != hash3


@pytest.mark.asyncio
async def test_check_deduplication():
    mock_pool = MagicMock()
    mock_conn = AsyncMock()
    mock_pool.acquire.return_value.__aenter__.return_value = mock_conn

    existing_doc_id = uuid.uuid4()
    mock_conn.fetchrow.return_value = {
        "id": existing_doc_id,
        "filename": "existing.pdf",
        "source_type": "pdf",
        "content_hash": "sample_hash",
        "status": "complete",
        "chunk_count": 10,
        "metadata": "{}",
        "created_at": None,
    }

    result = await check_deduplication("sample_hash", pool=mock_pool)
    assert result is not None
    assert result["id"] == existing_doc_id
    assert result["status"] == "complete"


# --- 2. Extractor Tests ---

def test_markdown_table_formatting():
    table = [
        ["Model", "Provider", "Cost"],
        ["gpt-4o", "OpenAI", "$2.50"],
        ["claude-3-5", "Anthropic", "$3.00"],
    ]
    md = format_markdown_table(table)
    assert "| Model | Provider | Cost |" in md
    assert "| --- | --- | --- |" in md
    assert "| gpt-4o | OpenAI | $2.50 |" in md


@pytest.mark.asyncio
async def test_pdf_extractor_digital_and_scanned():
    extractor = PDFExtractor()

    # Mock pypdfium2 doc with 2 pages: 1 digital, 1 scanned (<50 chars)
    mock_doc = MagicMock()
    mock_page1 = MagicMock()
    mock_page1.get_textpage.return_value.get_text_range.return_value = "This is a digital page with plenty of text that exceeds the fifty characters threshold for scanning detection."
    
    mock_page2 = MagicMock()
    mock_page2.get_textpage.return_value.get_text_range.return_value = "Scan"  # < 50 chars
    mock_pil_img = Image.new("RGB", (100, 100), color="white")
    mock_page2.render.return_value.to_pil.return_value = mock_pil_img

    mock_doc.__len__.return_value = 2
    mock_doc.__getitem__.side_effect = [mock_page1, mock_page2]

    # Mock pdfplumber for tables
    mock_plumber_doc = MagicMock()
    mock_plumb_p1 = MagicMock(extract_tables=MagicMock(return_value=[[["Col1", "Col2"], ["Val1", "Val2"]]]))
    mock_plumb_p2 = MagicMock(extract_tables=MagicMock(return_value=[]))
    mock_plumber_doc.pages = [mock_plumb_p1, mock_plumb_p2]

    with patch("pypdfium2.PdfDocument", return_value=mock_doc), \
         patch("pdfplumber.open", return_value=MagicMock(__enter__=MagicMock(return_value=mock_plumber_doc))), \
         patch("pytesseract.image_to_string", return_value="OCR Extracted Text from scanned document"):
        
        pages = await extractor.extract(b"%PDF-1.4 dummy bytes")

        assert len(pages) == 3  # Page 1 text, Page 1 table, Page 2 OCR text
        
        p1_text = [p for p in pages if p.page_number == 1 and p.content_type == "text"][0]
        assert "digital page" in p1_text.content
        assert p1_text.metadata["is_scanned"] is False

        p1_table = [p for p in pages if p.page_number == 1 and p.content_type == "table"][0]
        assert "| Col1 | Col2 |" in p1_table.content

        p2_ocr = [p for p in pages if p.page_number == 2 and p.content_type == "text"][0]
        assert "OCR Extracted Text" in p2_ocr.content
        assert p2_ocr.metadata["is_scanned"] is True


@pytest.mark.asyncio
async def test_docx_extractor():
    extractor = DocxExtractor()

    # Create in-memory DOCX
    doc = docx.Document()
    doc.add_heading("Architecture Overview", level=1)
    doc.add_paragraph("This is the intro paragraph.")
    doc.add_heading("Sub-section Details", level=2)
    doc.add_paragraph("Details paragraph under H2.")

    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Header 1"
    table.cell(0, 1).text = "Header 2"
    table.cell(1, 0).text = "A"
    table.cell(1, 1).text = "B"

    buf = io.BytesIO()
    doc.save(buf)
    docx_bytes = buf.getvalue()

    pages = await extractor.extract(docx_bytes)

    assert len(pages) >= 4
    # Check H1 heading hierarchy preservation
    h1_page = [p for p in pages if p.metadata.get("level") == "h1"][0]
    assert h1_page.content == "Architecture Overview"
    assert h1_page.metadata["section"] == "Architecture Overview"

    # Check paragraph under H1 inherits section
    intro_page = [p for p in pages if p.content == "This is the intro paragraph."][0]
    assert intro_page.metadata["section"] == "Architecture Overview"

    # Check table extraction
    table_page = [p for p in pages if p.content_type == "table"][0]
    assert "| Header 1 | Header 2 |" in table_page.content


@pytest.mark.asyncio
async def test_image_extractor():
    extractor = ImageExtractor()

    # Create 1200x800 test image (should be resized to max 1024)
    img = Image.new("RGB", (1200, 800), color="blue")
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    img_bytes = buf.getvalue()

    mock_client = MagicMock(spec=NeuroFlowClient)
    mock_client.chat = AsyncMock(
        return_value=GenerationResult(
            content="A detailed blue technical diagram with 3 boxes.",
            model="gpt-4o",
            input_tokens=100,
            output_tokens=50,
            latency_ms=150.0,
            cost_usd=0.001,
            finish_reason="stop",
        )
    )

    with patch("pytesseract.image_to_string", return_value="Diagram Header 2026"):
        pages = await extractor.extract(img_bytes, client=mock_client)

        assert len(pages) == 1
        page = pages[0]
        assert page.content_type == "image_description"
        assert "A detailed blue technical diagram" in page.content
        assert "Text found in image: Diagram Header 2026" in page.content
        assert page.metadata["original_size"] == (1200, 800)
        assert max(page.metadata["processed_size"]) <= 1024


@pytest.mark.asyncio
async def test_csv_extractor_small_and_large():
    extractor = CSVExtractor()

    # 1. Small CSV (150 rows -> 2 pages)
    small_df = pd.DataFrame({
        "id": range(150),
        "name": [f"Item-{i}" for i in range(150)],
        "price": [10.5 * i for i in range(150)],
    })
    small_buf = io.BytesIO()
    small_df.to_csv(small_buf, index=False)

    pages_small = await extractor.extract(small_buf.getvalue())
    assert len(pages_small) == 2
    assert pages_small[0].content_type == "table"
    assert pages_small[1].content_type == "table"
    assert pages_small[0].metadata["row_start"] == 0
    assert pages_small[0].metadata["row_end"] == 100
    assert pages_small[1].metadata["row_start"] == 100
    assert pages_small[1].metadata["row_end"] == 150

    # 2. Large CSV (1050 rows -> 1 statistical summary page + 11 table block pages = 12 pages)
    large_df = pd.DataFrame({
        "metric_a": range(1050),
        "category": ["Alpha" if i % 2 == 0 else "Beta" for i in range(1050)],
    })
    large_buf = io.BytesIO()
    large_df.to_csv(large_buf, index=False)

    pages_large = await extractor.extract(large_buf.getvalue())
    assert len(pages_large) == 12
    summary_page = pages_large[0]
    assert summary_page.metadata.get("is_statistical_summary") is True
    assert "CSV Dataset Summary" in summary_page.content
    assert "Numeric Summary" in summary_page.content


@pytest.mark.asyncio
async def test_url_extractor():
    extractor = URLExtractor()

    mock_resp = MagicMock()
    mock_resp.status_code = 200
    mock_resp.text = "<html><head><title>NeuroFlow Doc</title></head><body><h1>Welcome to NeuroFlow</h1><p>Main content text here.</p></body></html>"

    with patch("httpx.AsyncClient.get", AsyncMock(return_value=mock_resp)), \
         patch.object(extractor, "_check_robots_permission", AsyncMock(return_value=True)):

        pages = await extractor.extract("https://example.com/docs")
        assert len(pages) == 1
        page = pages[0]
        assert "Welcome to NeuroFlow" in page.content
        assert page.metadata["url"] == "https://example.com/docs"


@pytest.mark.asyncio
async def test_pptx_extractor():
    extractor = PPTXExtractor()

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "NeuroFlow AI Slide"
    slide.placeholders[1].text = "Bullet 1: Architecture\nBullet 2: Ingestion"
    
    # Speaker note
    slide.notes_slide.notes_text_frame.text = "Remember to explain pgvector."

    buf = io.BytesIO()
    prs.save(buf)

    pages = await extractor.extract(buf.getvalue())
    assert len(pages) == 1
    assert "NeuroFlow AI Slide" in pages[0].content
    assert "Bullet 1: Architecture" in pages[0].content
    assert "Remember to explain pgvector." in pages[0].content
    assert pages[0].metadata["has_speaker_notes"] is True


# --- 3. Chunker Strategy Tests ---

def test_fixed_size_chunking_sentence_boundaries():
    # Long text with clear sentence breaks
    sentences = [
        f"Sentence number {i} is full of descriptive information about vector databases and retrieval."
        for i in range(50)
    ]
    text = " ".join(sentences)

    chunks = fixed_size_chunking(text, target_tokens=100, overlap_tokens=20)
    assert len(chunks) > 1

    for chunk in chunks:
        # Assert sentence boundary: chunk ends with period
        assert chunk.content.endswith(".")
        assert chunk.token_count > 0


@pytest.mark.asyncio
async def test_semantic_chunking_topic_drops():
    text = "Machine learning algorithms learn patterns from training data. Deep neural networks use backpropagation. " \
           "The Eiffel Tower is located in Paris France. Tourists visit the Louvre museum to see famous art."

    mock_client = MagicMock(spec=NeuroFlowClient)
    # Return dissimilar embeddings between topic 1 (ML) and topic 2 (Paris)
    mock_client.embed = AsyncMock(
        return_value=[
            [1.0, 0.0, 0.0],  # ML 1
            [0.9, 0.1, 0.0],  # ML 2
            [0.0, 1.0, 0.0],  # Paris 1 (sim = 0.0 < 0.7)
            [0.1, 0.9, 0.0],  # Paris 2
        ]
    )

    chunks = await semantic_chunking(text, similarity_threshold=0.7, client=mock_client)
    assert len(chunks) == 2
    assert "Machine learning" in chunks[0].content
    assert "Eiffel Tower" in chunks[1].content


def test_hierarchical_chunking():
    pages = [
        ExtractedPage(page_number=1, content="Introduction to RAG", content_type="text", metadata={"level": "h1", "section": "Intro"}),
        ExtractedPage(page_number=1, content="RAG combines retrieval with LLM generation.", content_type="text", metadata={"level": "h2", "section": "Intro"}),
        ExtractedPage(page_number=1, content="Advanced Vector Stores", content_type="text", metadata={"level": "h1", "section": "Vectors"}),
        ExtractedPage(page_number=1, content="pgvector enables efficient HNSW cosine indexing.", content_type="text", metadata={"level": "h2", "section": "Vectors"}),
    ]

    chunks = hierarchical_chunking(pages)
    assert len(chunks) == 4

    parent_chunks = [c for c in chunks if c.metadata.get("is_parent")]
    child_chunks = [c for c in chunks if c.metadata.get("is_child")]

    assert len(parent_chunks) == 2
    assert len(child_chunks) == 2
    # Verify child has parent_id reference
    assert child_chunks[0].metadata["parent_id"] == parent_chunks[0].metadata["chunk_id"]
    assert child_chunks[1].metadata["parent_id"] == parent_chunks[1].metadata["chunk_id"]


def test_chunking_strategy_auto_selection():
    # 1. Table content -> fixed_size
    table_pages = [ExtractedPage(page_number=1, content="| A | B |", content_type="table")]
    assert select_chunking_strategy(table_pages, source_type="pdf") == "fixed_size"
    assert select_chunking_strategy([], source_type="csv") == "fixed_size"

    # 2. DOCX with headings -> hierarchical
    docx_pages = [ExtractedPage(page_number=1, content="Header", content_type="text", metadata={"level": "h1"})]
    assert select_chunking_strategy(docx_pages, source_type="docx") == "hierarchical"

    # 3. PDF > 50 pages -> semantic
    pdf_pages = [ExtractedPage(page_number=1, content="Long text", content_type="text", metadata={"total_pages": 60})]
    assert select_chunking_strategy(pdf_pages, source_type="pdf") == "semantic"

    # 4. Default -> fixed_size
    default_pages = [ExtractedPage(page_number=1, content="Short text", content_type="text")]
    assert select_chunking_strategy(default_pages, source_type="text") == "fixed_size"


# --- 4. Pipeline & OpenTelemetry Observability Tests ---

@pytest.mark.asyncio
async def test_ingestion_pipeline_end_to_end():
    mock_client = MagicMock(spec=NeuroFlowClient)
    mock_client.embed = AsyncMock(return_value=[[0.1] * 1536, [0.2] * 1536])

    mock_pool = MagicMock()
    mock_conn = AsyncMock()
    mock_tx = MagicMock()
    mock_tx.__aenter__ = AsyncMock(return_value=None)
    mock_tx.__aexit__ = AsyncMock(return_value=None)
    mock_conn.transaction = MagicMock(return_value=mock_tx)
    mock_pool.acquire.return_value.__aenter__.return_value = mock_conn

    pipeline = IngestionPipeline(client=mock_client, pool=mock_pool)

    doc_id = str(uuid.uuid4())
    sample_text = "NeuroFlow provides automated multi-modal ingestion. It extracts chunks and stores embeddings."

    with patch.object(pipeline.extractors["pdf"], "extract", AsyncMock(return_value=[
        ExtractedPage(page_number=1, content=sample_text, content_type="text", metadata={"total_pages": 1})
    ])), patch("backend.pipelines.ingestion.pipeline.tracer") as mock_tracer:
        
        mock_span = MagicMock()
        mock_tracer.start_as_current_span.return_value.__enter__.return_value = mock_span

        result = await pipeline.process_document(
            document_id=doc_id,
            source=b"dummy pdf bytes",
            source_type="pdf",
            filename="sample.pdf",
        )

        assert result["status"] == "complete"
        assert result["page_count"] == 1
        assert result["chunk_count"] >= 1

        # Verify OpenTelemetry span attributes
        mock_span.set_attribute.assert_any_call("document_id", doc_id)
        mock_span.set_attribute.assert_any_call("source_type", "pdf")
        mock_span.set_attribute.assert_any_call("page_count", 1)


# --- 5. API Ingestion Endpoints Tests ---

@pytest.mark.asyncio
async def test_api_ingest_and_deduplication():
    from fastapi.testclient import TestClient
    from backend.main import app

    client = TestClient(app)

    # 1. Test POST /ingest upload
    with patch("backend.api.ingest.get_pool", return_value=None), \
         patch("backend.api.ingest.enqueue_ingest_job", AsyncMock()) as mock_enqueue:

        file_content = b"Unique content for new ingestion document test."
        response = client.post(
            "/ingest",
            files={"file": ("test_doc.pdf", io.BytesIO(file_content), "application/pdf")},
        )
        assert response.status_code == 202
        data = response.json()
        assert data["status"] == "queued"
        assert data["duplicate"] is False
        assert "document_id" in data
        mock_enqueue.assert_awaited_once()

    # 2. Test Deduplication: second upload of same content returns duplicate=True
    with patch("backend.api.ingest.get_pool", return_value=None), \
         patch("backend.api.ingest.check_deduplication", AsyncMock(return_value={"id": "existing-uuid-1234", "status": "complete"})):

        response_dup = client.post(
            "/ingest",
            files={"file": ("test_doc.pdf", io.BytesIO(file_content), "application/pdf")},
        )
        assert response_dup.status_code == 202
        data_dup = response_dup.json()
        assert data_dup["document_id"] == "existing-uuid-1234"
        assert data_dup["duplicate"] is True
        assert data_dup["status"] == "complete"
