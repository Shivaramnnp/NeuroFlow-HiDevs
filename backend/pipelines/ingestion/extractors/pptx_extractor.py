from __future__ import annotations

import io
import logging
from typing import Any, List, Optional, Union

from pptx import Presentation

from .base import BaseExtractor, ExtractedPage
try:
    from backend.providers.client import NeuroFlowClient
except ImportError:
    from providers.client import NeuroFlowClient

logger = logging.getLogger("neuroflow-pptx-extractor")


class PPTXExtractor(BaseExtractor):
    """
    PowerPoint (PPTX) Presentation Extractor:
    - Extracts slide shapes, text frames, titles, and speaker notes
    - Identifies visual elements/diagrams
    - Each slide becomes one ExtractedPage
    """

    def __init__(self, client: Optional[NeuroFlowClient] = None):
        self.client = client

    async def extract(self, source: Union[str, bytes, io.BytesIO], **kwargs) -> List[ExtractedPage]:
        if isinstance(source, bytes):
            pptx_file = io.BytesIO(source)
        elif isinstance(source, io.BytesIO):
            pptx_file = source
        else:
            pptx_file = source

        prs = Presentation(pptx_file)
        results: List[ExtractedPage] = []

        for slide_idx, slide in enumerate(prs.slides):
            slide_num = slide_idx + 1
            text_blocks: List[str] = []
            has_images_or_shapes = False

            # Extract slide title and text frames
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        text = p.text.strip()
                        if text:
                            text_blocks.append(text)
                elif shape.shape_type in [13, 6, 7]:  # Picture, Group, Embedded OLE
                    has_images_or_shapes = True

            # Extract speaker notes if available
            notes_text = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    notes_text = notes

            content_lines = []
            if text_blocks:
                content_lines.append("\n".join(text_blocks))
            if notes_text:
                content_lines.append(f"\nSpeaker Notes:\n{notes_text}")

            full_content = "\n\n".join(content_lines).strip()
            if not full_content:
                full_content = f"[Slide {slide_num} - Visual/Diagram Slide]"

            results.append(
                ExtractedPage(
                    page_number=slide_num,
                    content=full_content,
                    content_type="text",
                    metadata={
                        "slide_number": slide_num,
                        "has_visual_elements": has_images_or_shapes,
                        "has_speaker_notes": bool(notes_text),
                    },
                )
            )

        return results
